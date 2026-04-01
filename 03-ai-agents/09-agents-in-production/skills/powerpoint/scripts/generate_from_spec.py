"""
Generate a PowerPoint presentation from a master spec JSON file.
This script creates text-based slides without requiring pre-generated images.

Usage:
    python generate_from_spec.py --spec-file /path/to/presentation_master_spec.json \
                                  --output-file /path/to/output.pptx \
                                  [--style-file /path/to/visual_style_guide.json]
"""

import argparse
import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def hex_to_rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.lstrip("#")
    r, g, b = int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
    return RGBColor(r, g, b)


def add_text_box(slide, text, left, top, width, height, font_size, bold=False, color="#FFFFFF", align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = hex_to_rgb(color)
    return txBox


def fill_background(slide, color: str):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(color)


def add_bullet_slide(prs, slide_data: dict, style: dict):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    bg_color = style.get("primary_color", "#1B2A4A")
    accent = style.get("accent_color", "#00C2CB")
    text_color = style.get("text_color", "#FFFFFF")

    fill_background(slide, bg_color)

    slide_type = slide_data.get("type", "content")
    title = slide_data.get("title", "")
    subtitle = slide_data.get("subtitle", "")
    bullets = slide_data.get("bullets", [])
    objective = slide_data.get("objective", "")
    main_message = slide_data.get("main_message", "")

    if slide_type == "capa":
        # Large centered title
        add_text_box(slide, title, 0.5, 2.0, 12.3, 1.5, 40, bold=True, color=text_color, align=PP_ALIGN.CENTER)
        if subtitle:
            add_text_box(slide, subtitle, 0.5, 3.8, 12.3, 0.8, 20, bold=False, color=accent, align=PP_ALIGN.CENTER)
        course = "Jornada de Dados — Engenheiro de IA"
        add_text_box(slide, course, 0.5, 5.0, 12.3, 0.6, 16, bold=False, color=accent, align=PP_ALIGN.CENTER)

    elif slide_type == "sobre_mim":
        add_text_box(slide, title, 0.5, 0.3, 12.3, 0.7, 28, bold=True, color=accent)
        bio = slide_data.get("bio", "")
        if bio:
            add_text_box(slide, bio, 0.5, 1.2, 12.3, 0.6, 16, color=text_color)
        y = 2.1
        for bullet in bullets:
            add_text_box(slide, f"• {bullet}", 0.7, y, 11.5, 0.5, 15, color=text_color)
            y += 0.55
        if main_message:
            add_text_box(slide, main_message, 0.5, 6.5, 12.3, 0.6, 14, color=accent)

    elif slide_type == "bloco_intro":
        # Accent bar on left
        bar = slide.shapes.add_shape(1, Inches(0.3), Inches(0.5), Inches(0.1), Inches(6.5))
        bar.fill.solid()
        bar.fill.fore_color.rgb = hex_to_rgb(accent)
        bar.line.fill.background()

        add_text_box(slide, title, 0.6, 0.5, 12.0, 0.9, 32, bold=True, color=accent)
        if objective:
            add_text_box(slide, objective, 0.6, 1.6, 12.0, 0.6, 16, color=text_color)
        y = 2.5
        for bullet in bullets:
            add_text_box(slide, f"→  {bullet}", 0.8, y, 11.5, 0.5, 15, color=text_color)
            y += 0.6
        if main_message:
            add_text_box(slide, main_message, 0.6, 6.6, 12.0, 0.55, 13, color=accent)

    elif slide_type == "conclusao":
        add_text_box(slide, title, 0.5, 0.3, 12.3, 0.9, 30, bold=True, color=accent, align=PP_ALIGN.CENTER)
        y = 1.5
        for bullet in bullets:
            add_text_box(slide, f"✓  {bullet}", 0.8, y, 11.5, 0.55, 16, color=text_color)
            y += 0.65
        if main_message:
            add_text_box(slide, main_message, 0.5, 6.5, 12.3, 0.6, 14, color=accent, align=PP_ALIGN.CENTER)

    else:  # content
        # Thin accent top bar
        bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.07))
        bar.fill.solid()
        bar.fill.fore_color.rgb = hex_to_rgb(accent)
        bar.line.fill.background()

        add_text_box(slide, title, 0.5, 0.25, 12.3, 0.8, 26, bold=True, color=accent)
        if objective:
            add_text_box(slide, objective, 0.5, 1.15, 12.3, 0.55, 15, color=text_color)
        y = 1.85 if objective else 1.3
        for bullet in bullets:
            add_text_box(slide, f"•  {bullet}", 0.7, y, 11.5, 0.55, 15, color=text_color)
            y += 0.6
        if main_message:
            add_text_box(slide, main_message, 0.5, 6.6, 12.3, 0.55, 13, color=accent)

    # Speaker notes
    notes_parts = []
    if objective:
        notes_parts.append(f"Objetivo: {objective}")
    if bullets:
        notes_parts.append("Bullets:\n" + "\n".join(f"  • {b}" for b in bullets))
    if main_message:
        notes_parts.append(f"Mensagem principal: {main_message}")
    if notes_parts:
        slide.notes_slide.notes_text_frame.text = "\n\n".join(notes_parts)

    return slide


def generate_ppt_from_spec(spec_file: str, output_file: str, style_file: str | None = None) -> str:
    with open(spec_file, "r", encoding="utf-8") as f:
        spec = json.load(f)

    style = {}
    if style_file and Path(style_file).exists():
        with open(style_file, "r", encoding="utf-8") as f:
            style = json.load(f)
    elif spec.get("brand"):
        # Try to load from spec's brand reference
        brand_ref = spec["brand"].get("source", "")
        if brand_ref:
            candidate = Path(spec_file).parent / brand_ref
            if candidate.exists():
                with open(candidate, "r", encoding="utf-8") as f:
                    style = json.load(f)

    # Fallback style
    style.setdefault("primary_color", "#1B2A4A")
    style.setdefault("accent_color", "#00C2CB")
    style.setdefault("text_color", "#FFFFFF")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slides = spec.get("slides", [])
    if not slides:
        return "Error: No slides found in spec file."

    for slide_data in slides:
        add_bullet_slide(prs, slide_data, style)

    Path(output_file).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_file)
    return f"Successfully generated {len(slides)}-slide presentation at {output_file}"


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Generate PPTX from master spec JSON")
    parser.add_argument("--spec-file", required=True, help="Path to presentation_master_spec.json")
    parser.add_argument("--output-file", required=True, help="Output path for .pptx file")
    parser.add_argument("--style-file", default=None, help="Optional path to visual_style_guide.json")
    args = parser.parse_args()

    try:
        print(generate_ppt_from_spec(args.spec_file, args.output_file, args.style_file))
    except Exception as e:
        print(f"Error: {e}")
        raise
